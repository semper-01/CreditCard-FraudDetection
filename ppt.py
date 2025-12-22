# You may need to run: pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_slide(title, points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Credit Card Fraud Detection System"
slide.placeholders[1].text = "Authors: Leonce M. & Grace U.\nCollege of Science and Technology\nML Final Project"

# Slide 2: Project Goals
add_slide("Project Introduction & Goals", [
    "Problem: High-stakes Binary Classification on imbalanced data (0.17% fraud).",
    "Goal: Maximize detection (Recall) while minimizing false alerts (Precision).",
    "Key Algorithm: XGBoost with SHAP explainability."
])

# Slide 3: ML Project Cycle
add_slide("The ML Project Cycle", [
    "1. Define: Fraud detection problem identification.",
    "2. Understand: Data analysis and feature distributions.",
    "3. Preprocess: Feature engineering (AmountLog, TimeHours) and scaling.",
    "4. Manage: Model training (XGBoost) and Evaluation.",
    "5. Deploy: Global accessibility via FastAPI and ASP.NET."
])

# Slide 4: Data Preprocessing
add_slide("Data Preparation (preprocess.py)", [
    "Feature Engineering: Created AmountLog and Time-based features.",
    "Handling Imbalance: Applied SMOTE to increase minority class representation.",
    "Scaling: Robust scaling using StandardScaler saved as scaler.pkl.",
    "Split: Time-series split to ensure real-world simulation."
])

# Slide 5: Model & Optimization
add_slide("XGBoost & Threshold Tuning", [
    "Algorithm: XGBoost with scale_pos_weight for imbalance.",
    "Optimization: Iterative threshold testing (threshold.py).",
    "Best Threshold: 0.931 (93.1% confidence required).",
    "Result: F1-Score improved from 0.74 to 0.81."
])

# Slide 6: MLOps Architecture
add_slide("3-Tier MLOps Deployment", [
    "Presentation Layer: ASP.NET Core UI.",
    "Service Layer: FastAPI REST API.",
    "Model Layer: Serialized artifacts (model.pkl, threshold.json).",
    "Explainability: Real-time SHAP values via the /explain endpoint."
])

prs.save('Fraud_Detection_Presentation.pptx')
print("PowerPoint generated successfully!")